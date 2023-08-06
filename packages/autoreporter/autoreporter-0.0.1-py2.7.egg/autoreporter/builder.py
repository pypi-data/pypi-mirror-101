# https://stackoverflow.com/questions/18897511/how-to-drawimage-a-matplotlib-figure-in-a-reportlab-canvas

def fig2img(fig):
    imgdata = cStringIO.StringIO()
    fig.savefig(imgdata, format='png')
    imgdata.seek(0)  # rewind the data

    return ImageReader(imgdata)

def builder(path, ref, fig_dict):
    SCALE = 2
    PAD = 0.3
    from pptx import Presentation
    from pptx.util import Inches
    import tempfile
    
    prs = Presentation(ref)
    out = Presentation()

    out.slide_layouts = prs.slide_layouts
    
    out.slide_width = prs.slide_width
    out.slide_height = prs.slide_height
    
    for slide in prs.slides:
        blank_slide_layout = out.slide_layouts[6]
        o_slide = out.slides.add_slide(blank_slide_layout)
        
        for shape in slide.shapes:
            if shape.text[0] == "!" and shape.text[1] == "(" and shape.text[-1] == ")":
                f = tempfile.NamedTemporaryFile(suffix=".png")
                fig = fig_dict[shape.text[2:-1]]
                fig.set_size_inches(SCALE*shape.width/Inches(1), SCALE*shape.height/Inches(1))
                fig.set_tight_layout({"pad" : PAD})
                fig.savefig(f.name, format='png',dpi = 300)
                pic = o_slide.shapes.add_picture(f.name, shape.left, shape.top, width=shape.width)
            else:
                # o_slide.shapes.add_shape(shape.shape_type, shape.left, shape.top, shape.width, shape.height)
                # o_slide.shapes[-1] = shape
                o_slide.append(shape)
                
    out.save(path)
